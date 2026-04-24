"""Render a PPTX to per-slide PNGs via soffice → PDF → pdftoppm, and run simple
text-based QA: flag untouched placeholders and verify shape-count parity with
the source template.
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

UNTOUCHED_MARKERS = (
    "메인 타이틀을 입력하세요",
    "서브 타이틀을 입력하세요",
    "타이틀을 입력하세요",
    "타이틀 입력",
    "목차를 입력해주세요",
    "내용을 입력하세요",
    "내용 입력",
)

SECTION_EXEMPT = ("PRESENTATION TEMPLATE", "LOGO HERE")

SOFFICE = "/opt/homebrew/bin/soffice"


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def render(pptx: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir = out_dir / "_pdf"
    pdf_dir.mkdir(exist_ok=True)
    subprocess.run(
        [SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(pptx)],
        check=True,
        capture_output=True,
    )
    pdfs = list(pdf_dir.glob("*.pdf"))
    if not pdfs:
        raise RuntimeError("soffice produced no PDF")
    pdf = pdfs[0]
    prefix = out_dir / "slide"
    subprocess.run(
        ["pdftoppm", "-png", "-r", "150", str(pdf), str(prefix)], check=True, capture_output=True
    )
    shutil.rmtree(pdf_dir, ignore_errors=True)
    return sorted(out_dir.glob("slide-*.png"))


def qa(pptx: Path, classification_path: Path | None, plan_path: Path | None) -> list[dict]:
    prs = Presentation(str(pptx))
    issues: list[dict] = []

    classification = {}
    if classification_path and classification_path.exists():
        data = json.loads(classification_path.read_text())
        classification = {s["index"]: s for s in data["slides"]}

    plan = None
    if plan_path and plan_path.exists():
        plan = json.loads(plan_path.read_text())["slides"]

    for position, slide in enumerate(prs.slides, start=1):
        shape_count = sum(1 for _ in iter_shapes(slide.shapes))
        text_full = "\n".join(
            shape.text_frame.text for shape in iter_shapes(slide.shapes) if shape.has_text_frame
        )

        plan_entry = plan[position - 1] if plan and position - 1 < len(plan) else None
        kind = plan_entry["kind"] if plan_entry else "?"
        source_index = plan_entry["source_index"] if plan_entry else None

        for marker in UNTOUCHED_MARKERS:
            if marker in text_full:
                if kind == "section" and any(ex in text_full for ex in SECTION_EXEMPT):
                    continue
                issues.append(
                    {
                        "position": position,
                        "source_index": source_index,
                        "kind": kind,
                        "reason": f"untouched placeholder '{marker}'",
                    }
                )
                break

        if source_index is not None and source_index in classification:
            expected = classification[source_index]["shape_count"]
            if shape_count != expected:
                issues.append(
                    {
                        "position": position,
                        "source_index": source_index,
                        "kind": kind,
                        "reason": f"shape count {shape_count} != template {expected}",
                    }
                )
    return issues


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out-dir", required=True)
    ap.add_argument("--classification", default=None)
    ap.add_argument("--plan", default=None)
    args = ap.parse_args()

    pptx = Path(args.pptx).expanduser().resolve()
    out_dir = Path(args.out_dir).expanduser().resolve()

    if not Path(SOFFICE).exists():
        sys.exit(f"soffice not found at {SOFFICE}")

    pngs = render(pptx, out_dir)
    issues = qa(
        pptx,
        Path(args.classification) if args.classification else None,
        Path(args.plan) if args.plan else None,
    )

    report = {
        "output_path": str(pptx),
        "slide_count": len(pngs),
        "preview_dir": str(out_dir),
        "issues": issues,
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
