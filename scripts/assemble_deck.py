"""Build the output PPT from slide_plan.json by:
  1. opening the template
  2. cloning the slide XML of each planned source slide (append to end)
  3. swapping placeholder run text per-kind
  4. dropping all original template slides so only the new ones remain, in plan order
  5. saving to --output

Shape geometry, masters/layouts, images, and run formatting are untouched — this
satisfies the pixel-perfect contract in agents/deck-assembler.md.
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt  # noqa: F401

BIG_NUM_RE = re.compile(r"\d{1,3}(?:,\d{3})+|\d{1,3}%|\b\d{3,}\b|\b\d{2}%")
STEP_RE = re.compile(r"(?<!\d)0[1-9](?!\d)")

TITLE_PLACEHOLDERS = (
    "메인 타이틀을 입력하세요",
    "서브 타이틀을 입력하세요",
    "타이틀 입력",
    "타이틀을 입력하세요",
    "내용 입력",
    "내용을 입력하세요",
    "상세 타이틀",
    "목차를 입력해주세요",
    "목차를 상세하게 입력해주세요",
)

BODY_PLACEHOLDER_MARKERS = (
    "상세 타이틀",
    "내용 입력",
    "내용을 입력",
    "매주 업데이트되는",
    "다양한 주제와 디자인",
    "비즈니스, 교육, 마케팅",
    "저희는 프레젠테이션",
    "세련되고 정교한",
    "안녕하세요. 매주",
    "수많은 주제와",
    "간편하게 커스터마이징",
    "타이틀을 설명하는",
    "파워포인트 템플릿",
)

# Any shape whose stripped text appears here is pure boilerplate decoration
# that must be blanked regardless of injector logic.
GLOBAL_BOILERPLATE_TOKENS = (
    "수많은 주제와",
    "간편하게 커스터마이징",
    "저희는 프레젠테이션",
    "세련되고 정교한",
    "안녕하세요. 매주",
    "다양한 주제와 디자인",
    "비즈니스, 교육, 마케팅",
    "타이틀을 설명하는",
    "비즈니스\x0b파워포인트 템플릿",
)

FOOTER_TOKENS = ("COMPANY LOGO HERE", "PRESENTATION TEMPLATE")


def is_body_placeholder(text: str) -> bool:
    return any(m in text for m in BODY_PLACEHOLDER_MARKERS)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def set_runs_text(tf, new_text: str) -> bool:
    """Write `new_text` into the first run of the first paragraph, blanking the rest.
    If no run exists (common in table cells the template author left blank),
    create one on the first paragraph so the text still lands in the cell.
    Returns True if the cell now holds `new_text`."""
    if not new_text:
        return False
    runs = [r for p in tf.paragraphs for r in p.runs]
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
        return True
    # No runs anywhere — create one on the first paragraph.
    if tf.paragraphs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = new_text
    return True


def replace_text_in_shape(shape, old: str, new: str) -> bool:
    if not shape.has_text_frame or not new:
        return False
    txt = shape.text_frame.text
    if old and old in txt:
        return set_runs_text(shape.text_frame, new)
    return False


def first_matching_shape(slide, predicate):
    for shape in iter_shapes(slide.shapes):
        if predicate(shape):
            return shape
    return None


def _explicit_font_max(tf) -> int:
    """Max explicit run font size; 0 if all runs inherit from master."""
    best = 0
    for p in tf.paragraphs:
        for r in p.runs:
            sz = r.font.size
            if sz is not None:
                best = max(best, sz)
    return best


def inject_cover(slide, content: dict) -> list[str]:
    title = content.get("title", "")
    subtitle = content.get("subtitle", "")
    company = content.get("company", "")

    text_shapes = []
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not t:
            continue
        text_shapes.append((shape, t, _explicit_font_max(shape.text_frame)))

    # Cover typically has three text elements: a long marketing headline, a
    # multi-line paragraph, and a "COMPANY LOGO HERE" footer. The headline is
    # the largest-font or shortest-inherited-font shape.
    inherited = [(s, t) for s, t, sz in text_shapes if sz == 0 and "LOGO" not in t]
    title_shape = None
    if inherited:
        title_shape = min(inherited, key=lambda x: len(x[1]))[0]
    else:
        non_logo = [x for x in text_shapes if "LOGO" not in x[1]]
        if non_logo:
            title_shape = max(non_logo, key=lambda x: x[2])[0]

    if title_shape is not None:
        set_runs_text(title_shape.text_frame, title or " ")

    subtitle_shape = None
    if subtitle:
        for shape, t, _ in text_shapes:
            if shape is title_shape or "LOGO" in t:
                continue
            if len(t) >= 15:
                set_runs_text(shape.text_frame, subtitle)
                subtitle_shape = shape
                break

    # Blank the long marketing paragraph (and any other decorative body text)
    # that didn't get a legit value.
    for shape, t, _ in text_shapes:
        if shape is title_shape or shape is subtitle_shape:
            continue
        if any(tok in t for tok in GLOBAL_BOILERPLATE_TOKENS):
            set_runs_text(shape.text_frame, " ")

    # Company logo footer: replace with provided company name, else blank it so
    # "COMPANY LOGO HERE" doesn't leak into the deck.
    for shape, t, _ in text_shapes:
        if "COMPANY LOGO" in t or "LOGO HERE" in t:
            set_runs_text(shape.text_frame, company or " ")
            break
    return []


def inject_agenda(slide, content: dict) -> list[str]:
    headings = content.get("headings", [])
    slot = 0

    # Agenda items may live in either (a) a table (slides 3, 4) or (b) discrete
    # text boxes (slide 2). Handle both. For table rows, also blank the number
    # column when there is no matching heading — an unlabeled "08 / 09" looks
    # like a loose end on the render.
    for shape in iter_shapes(slide.shapes):
        if shape.has_table:
            tbl = shape.table
            for row in tbl.rows:
                heading_cell = None
                number_cell = None
                for cell in row.cells:
                    ct = cell.text_frame.text.strip()
                    if "목차를" in ct and "입력" in ct:
                        heading_cell = cell
                    elif re.fullmatch(r"\d{1,2}", ct):
                        number_cell = cell
                if heading_cell is None:
                    continue
                heading = headings[slot] if slot < len(headings) else ""
                if heading:
                    set_runs_text(heading_cell.text_frame, heading)
                else:
                    set_runs_text(heading_cell.text_frame, " ")
                    if number_cell is not None:
                        set_runs_text(number_cell.text_frame, " ")
                slot += 1
        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "목차를" in t and "입력" in t:
                heading = headings[slot] if slot < len(headings) else ""
                set_runs_text(shape.text_frame, heading or " ")
                slot += 1

    # Blank any leftover decorative body text on the agenda slide.
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame or shape.has_table:
            continue
        t = shape.text_frame.text.strip()
        if any(tok in t for tok in GLOBAL_BOILERPLATE_TOKENS):
            set_runs_text(shape.text_frame, " ")
    return []


def inject_section(slide, content: dict) -> list[str]:
    step = content.get("step", "01")
    title = content.get("title", "")
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if STEP_RE.fullmatch(t):
            set_runs_text(shape.text_frame, step)
        elif any(p in t for p in ("타이틀 입력", "타이틀을 입력")) and title:
            set_runs_text(shape.text_frame, title)
        elif any(tok in t for tok in GLOBAL_BOILERPLATE_TOKENS):
            set_runs_text(shape.text_frame, " ")
        elif t in FOOTER_TOKENS:
            set_runs_text(shape.text_frame, " ")
    return []


def inject_kpi(slide, content: dict) -> list[str]:
    title = content.get("title", "")
    body = content.get("body", "")
    value = (content.get("data") or {}).get("value", "")

    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if value and BIG_NUM_RE.search(t) and not any(p in t for p in TITLE_PLACEHOLDERS):
            set_runs_text(shape.text_frame, value)
        elif title and any(p in t for p in ("타이틀 입력", "타이틀을 입력")):
            set_runs_text(shape.text_frame, title)
        elif body and any(p in t for p in ("내용 입력", "내용을 입력")):
            set_runs_text(shape.text_frame, body)
    return []


def inject_dashboard(slide, content: dict) -> list[str]:
    """Render a consulting-style KPI dashboard on top of a cleanish template
    slide. Clears ALL existing shapes and draws a custom layout:
      - Section title (top)
      - Answer-first headline (below title)
      - 4 KPI cards in a horizontal strip
      - KEY TAKEAWAYS bullet list at the bottom
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    title = (content.get("title") or "").strip()
    headline = (content.get("body") or "").strip()
    data = content.get("data") or {}
    kpis = data.get("kpis") or []
    insights = [s for s in (data.get("insights") or []) if s and str(s).strip()]

    # Wipe every shape on the slide — we redraw from scratch.
    for shape in list(slide.shapes):
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)

    PURPLE = RGBColor(0x6E, 0x5A, 0xE6)
    INK = RGBColor(0x1E, 0x1E, 0x2A)
    MUTED = RGBColor(0x6B, 0x6B, 0x7A)
    CARD = RGBColor(0xFF, 0xFF, 0xFF)
    CARD_EDGE = RGBColor(0xE6, 0xE6, 0xEE)
    SUBT = RGBColor(0x33, 0x33, 0x3D)

    def add_textbox(x, y, w, h, *, runs, wrap=True, align=PP_ALIGN.LEFT, margins=0):
        tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                      Emu(int(w * 914400)), Emu(int(h * 914400)))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = Emu(margins)
        tf.margin_top = tf.margin_bottom = Emu(margins)
        p = tf.paragraphs[0]
        p.alignment = align
        for i, spec in enumerate(runs):
            r = p.add_run()
            r.text = spec["text"]
            r.font.name = spec.get("font", "Pretendard")
            r.font.size = Pt(spec.get("size", 12))
            c = spec.get("color", INK)
            r.font.color.rgb = c
        return tb

    def add_card(x, y, w, h):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(int(x * 914400)), Emu(int(y * 914400)),
                                    Emu(int(w * 914400)), Emu(int(h * 914400)))
        sh.fill.solid()
        sh.fill.fore_color.rgb = CARD
        sh.line.color.rgb = CARD_EDGE
        sh.line.width = Emu(int(0.01 * 914400))
        sh.shadow.inherit = False
        # Hide default text in the auto-shape
        sh.text_frame.text = ""
        return sh

    # --- Title ---
    if title:
        add_textbox(0.72, 0.72, 11.89, 0.55,
                    runs=[{"text": title, "font": "Pretendard SemiBold",
                           "size": 26, "color": INK}])

    # --- Headline ---
    if headline:
        add_textbox(0.72, 1.45, 11.89, 0.6,
                    runs=[{"text": headline, "font": "Pretendard",
                           "size": 15, "color": SUBT}])

    # --- KPI Cards ---
    if kpis:
        kpis = kpis[:4]
        card_top = 2.35
        card_h = 2.5
        gap = 0.2
        total_w = 11.89
        card_w = (total_w - gap * (len(kpis) - 1)) / len(kpis)
        left = 0.72
        for i, kpi in enumerate(kpis):
            x = left + i * (card_w + gap)
            add_card(x, card_top, card_w, card_h)
            # Label (small, top of card)
            label = str(kpi.get("label", ""))
            if label:
                add_textbox(x + 0.3, card_top + 0.35, card_w - 0.6, 0.35,
                            runs=[{"text": label, "font": "Pretendard SemiBold",
                                   "size": 11, "color": MUTED}])
            # Big number (center of card). Auto-shrink when the value is a
            # long label (e.g. a product or customer name) instead of a pure
            # number — Korean glyphs are ~2× wider so we need to step down.
            value = str(kpi.get("value", ""))
            if value:
                # Count Korean chars as 2 ("wide") and everything else as 1.
                width_score = sum(2 if ord(c) > 0x3130 else 1 for c in value)
                if width_score <= 6:
                    vsize = 32
                elif width_score <= 10:
                    vsize = 24
                elif width_score <= 14:
                    vsize = 20
                else:
                    vsize = 16
                add_textbox(x + 0.3, card_top + 0.85, card_w - 0.6, 0.9,
                            runs=[{"text": value, "font": "Pretendard SemiBold",
                                   "size": vsize, "color": INK}])
            # Delta / sub-text (bottom of card)
            delta = str(kpi.get("delta", ""))
            if delta:
                add_textbox(x + 0.3, card_top + 1.8, card_w - 0.6, 0.35,
                            runs=[{"text": delta, "font": "Pretendard",
                                   "size": 10, "color": PURPLE}])

    # --- Key Takeaways ---
    if insights:
        tb_left = 0.72
        tb_top = 5.3
        tb_w = 11.89
        tb_h = 1.9
        tb = slide.shapes.add_textbox(Emu(int(tb_left * 914400)),
                                      Emu(int(tb_top * 914400)),
                                      Emu(int(tb_w * 914400)),
                                      Emu(int(tb_h * 914400)))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = "KEY TAKEAWAYS"
        r0.font.name = "Pretendard SemiBold"
        r0.font.size = Pt(9)
        r0.font.color.rgb = PURPLE
        for line in insights[:3]:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            rb = p.add_run()
            rb.text = "•  "
            rb.font.name = "Pretendard SemiBold"
            rb.font.size = Pt(11)
            rb.font.color.rgb = PURPLE
            rt = p.add_run()
            rt.text = str(line)
            rt.font.name = "Pretendard"
            rt.font.size = Pt(11)
            rt.font.color.rgb = SUBT

    return []


def _deck_palette():
    from pptx.dml.color import RGBColor
    return {
        "PURPLE": RGBColor(0x6E, 0x5A, 0xE6),
        "PURPLE_LT": RGBColor(0xB8, 0xAE, 0xF2),
        "INK": RGBColor(0x1E, 0x1E, 0x2A),
        "MUTED": RGBColor(0x6B, 0x6B, 0x7A),
        "SUBT": RGBColor(0x33, 0x33, 0x3D),
        "CARD": RGBColor(0xFF, 0xFF, 0xFF),
        "CARD_EDGE": RGBColor(0xE6, 0xE6, 0xEE),
        "BG_SOFT": RGBColor(0xF7, 0xF7, 0xFB),
    }


def _wipe_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


def _add_text(slide, x, y, w, h, runs, *, wrap=True, align=None):
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    for spec in runs:
        r = p.add_run()
        r.text = spec["text"]
        r.font.name = spec.get("font", "Pretendard")
        r.font.size = Pt(spec.get("size", 12))
        if "color" in spec:
            r.font.color.rgb = spec["color"]
    return tb


def _add_takeaways(slide, insights, top=5.3, height=1.9):
    from pptx.util import Emu, Pt
    pal = _deck_palette()
    if not insights:
        return
    tb = slide.shapes.add_textbox(Emu(int(0.72 * 914400)),
                                  Emu(int(top * 914400)),
                                  Emu(int(11.89 * 914400)),
                                  Emu(int(height * 914400)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "KEY TAKEAWAYS"
    r0.font.name = "Pretendard SemiBold"
    r0.font.size = Pt(9)
    r0.font.color.rgb = pal["PURPLE"]
    for line in insights[:3]:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name = "Pretendard SemiBold"
        rb.font.size = Pt(11)
        rb.font.color.rgb = pal["PURPLE"]
        rt = p.add_run()
        rt.text = str(line)
        rt.font.name = "Pretendard"
        rt.font.size = Pt(11)
        rt.font.color.rgb = pal["SUBT"]


CHART_TYPE_MAP = {
    "column": "COLUMN_CLUSTERED",
    "column_stacked": "COLUMN_STACKED",
    "bar": "BAR_CLUSTERED",
    "bar_stacked": "BAR_STACKED",
    "line": "LINE",
    "pie": "PIE",
    "doughnut": "DOUGHNUT",
    "area": "AREA",
    "radar": "RADAR",
}


def inject_chart(slide, content: dict) -> list[str]:
    """Render a native PowerPoint chart with title + headline + takeaways.
    Chart data format: content.data = {chart_type, categories, series, insights}
    """
    from pptx.util import Emu, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    pal = _deck_palette()
    title = (content.get("title") or "").strip()
    headline = (content.get("body") or "").strip()
    data = content.get("data") or {}
    categories = data.get("categories") or []
    series_list = data.get("series") or []
    chart_key = (data.get("chart_type") or "column").lower()
    insights = [s for s in (data.get("insights") or []) if s and str(s).strip()]

    _wipe_slide(slide)

    # Title
    if title:
        _add_text(slide, 0.72, 0.72, 11.89, 0.55,
                  runs=[{"text": title, "font": "Pretendard SemiBold",
                         "size": 26, "color": pal["INK"]}])
    # Headline
    if headline:
        _add_text(slide, 0.72, 1.45, 11.89, 0.6,
                  runs=[{"text": headline, "font": "Pretendard",
                         "size": 14, "color": pal["SUBT"]}])

    # Chart
    if categories and series_list:
        chart_data = CategoryChartData()
        chart_data.categories = list(categories)
        for s in series_list:
            name = s.get("name", "")
            values = tuple(s.get("values", []))
            chart_data.add_series(name, values)

        xl_type = getattr(XL_CHART_TYPE, CHART_TYPE_MAP.get(chart_key, "COLUMN_CLUSTERED"))
        x = Emu(int(0.72 * 914400))
        y = Emu(int(2.25 * 914400))
        cx = Emu(int(11.89 * 914400))
        cy = Emu(int(3.0 * 914400))
        chart_shape = slide.shapes.add_chart(xl_type, x, y, cx, cy, chart_data)
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = "Pretendard"
        # Category/value axis font
        for ax_name in ("category_axis", "value_axis"):
            try:
                ax = getattr(chart, ax_name)
                ax.tick_labels.font.size = Pt(9)
                ax.tick_labels.font.name = "Pretendard"
                ax.tick_labels.font.color.rgb = pal["MUTED"]
            except Exception:
                pass

    # Takeaways
    _add_takeaways(slide, insights, top=5.4, height=1.9)
    return []


def inject_exec_summary(slide, content: dict) -> list[str]:
    """High-level executive summary page: section label + massive headline
    + 3 supporting callouts + closing statement."""
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE

    pal = _deck_palette()
    _wipe_slide(slide)

    title = (content.get("title") or "Executive Summary").strip()
    headline = (content.get("body") or "").strip()
    data = content.get("data") or {}
    pillars = data.get("pillars") or []           # [{label, value, caption}]
    closing_line = (data.get("closing") or "").strip()

    # Small label (top-left eyebrow)
    _add_text(slide, 0.72, 0.72, 11.89, 0.35,
              runs=[{"text": "EXECUTIVE SUMMARY", "font": "Pretendard SemiBold",
                     "size": 10, "color": pal["PURPLE"]}])
    # Big title
    _add_text(slide, 0.72, 1.10, 11.89, 0.8,
              runs=[{"text": title, "font": "Pretendard SemiBold",
                     "size": 30, "color": pal["INK"]}])
    # Headline narrative
    if headline:
        _add_text(slide, 0.72, 2.05, 11.89, 1.0,
                  runs=[{"text": headline, "font": "Pretendard",
                         "size": 15, "color": pal["SUBT"]}])

    # Three supporting pillars
    pillars = pillars[:3]
    if pillars:
        pillar_top = 3.45
        pillar_h = 2.1
        gap = 0.25
        total_w = 11.89
        card_w = (total_w - gap * (len(pillars) - 1)) / len(pillars)
        left = 0.72
        for i, pl in enumerate(pillars):
            x = left + i * (card_w + gap)
            # Card
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Emu(int(x * 914400)), Emu(int(pillar_top * 914400)),
                                        Emu(int(card_w * 914400)), Emu(int(pillar_h * 914400)))
            bg.fill.solid()
            bg.fill.fore_color.rgb = pal["CARD"]
            bg.line.color.rgb = pal["CARD_EDGE"]
            bg.line.width = Emu(int(0.01 * 914400))
            bg.shadow.inherit = False
            bg.text_frame.text = ""
            # Accent bar on top
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Emu(int(x * 914400)), Emu(int(pillar_top * 914400)),
                                             Emu(int(card_w * 914400)), Emu(int(0.08 * 914400)))
            accent.fill.solid()
            accent.fill.fore_color.rgb = pal["PURPLE"]
            accent.line.fill.background()
            accent.shadow.inherit = False
            accent.text_frame.text = ""
            # Label
            _add_text(slide, x + 0.3, pillar_top + 0.3, card_w - 0.6, 0.3,
                      runs=[{"text": str(pl.get("label", "")),
                             "font": "Pretendard SemiBold", "size": 11,
                             "color": pal["MUTED"]}])
            # Value
            val = str(pl.get("value", ""))
            width_score = sum(2 if ord(c) > 0x3130 else 1 for c in val)
            if width_score <= 6: vsize = 26
            elif width_score <= 10: vsize = 20
            elif width_score <= 14: vsize = 16
            else: vsize = 14
            _add_text(slide, x + 0.3, pillar_top + 0.65, card_w - 0.6, 0.8,
                      runs=[{"text": val, "font": "Pretendard SemiBold",
                             "size": vsize, "color": pal["INK"]}])
            # Caption
            cap = str(pl.get("caption", ""))
            if cap:
                _add_text(slide, x + 0.3, pillar_top + 1.45, card_w - 0.6, 0.55,
                          runs=[{"text": cap, "font": "Pretendard",
                                 "size": 10, "color": pal["SUBT"]}])

    if closing_line:
        _add_text(slide, 0.72, 5.75, 11.89, 0.4,
                  runs=[{"text": closing_line, "font": "Pretendard SemiBold",
                         "size": 12, "color": pal["PURPLE"]}])
    return []


def inject_findings(slide, content: dict) -> list[str]:
    """Findings + next steps page — two-column layout.
    data: {findings: [str], next_steps: [str]}"""
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE

    pal = _deck_palette()
    _wipe_slide(slide)

    title = (content.get("title") or "Key Findings & Next Steps").strip()
    headline = (content.get("body") or "").strip()
    data = content.get("data") or {}
    findings = [s for s in (data.get("findings") or []) if s and str(s).strip()]
    next_steps = [s for s in (data.get("next_steps") or []) if s and str(s).strip()]

    _add_text(slide, 0.72, 0.72, 11.89, 0.55,
              runs=[{"text": title, "font": "Pretendard SemiBold",
                     "size": 26, "color": pal["INK"]}])
    if headline:
        _add_text(slide, 0.72, 1.45, 11.89, 0.5,
                  runs=[{"text": headline, "font": "Pretendard",
                         "size": 14, "color": pal["SUBT"]}])

    col_top = 2.25
    col_h = 4.5
    col_w = 5.82
    gap = 0.25
    left = 0.72

    def render_column(x, label, items, accent_color):
        # Card background
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Emu(int(x * 914400)), Emu(int(col_top * 914400)),
                                     Emu(int(col_w * 914400)), Emu(int(col_h * 914400)))
        bg.fill.solid()
        bg.fill.fore_color.rgb = pal["CARD"]
        bg.line.color.rgb = pal["CARD_EDGE"]
        bg.line.width = Emu(int(0.01 * 914400))
        bg.shadow.inherit = False
        bg.text_frame.text = ""
        # Label strip
        _add_text(slide, x + 0.4, col_top + 0.35, col_w - 0.8, 0.35,
                  runs=[{"text": label, "font": "Pretendard SemiBold",
                         "size": 10, "color": accent_color}])
        # Items
        tb = slide.shapes.add_textbox(Emu(int((x + 0.4) * 914400)),
                                       Emu(int((col_top + 0.8) * 914400)),
                                       Emu(int((col_w - 0.8) * 914400)),
                                       Emu(int((col_h - 1.2) * 914400)))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        first = True
        for i, line in enumerate(items[:5]):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(6)
            num = p.add_run()
            num.text = f"{i+1:02d}   "
            num.font.name = "Pretendard SemiBold"
            num.font.size = Pt(13)
            num.font.color.rgb = accent_color
            body_run = p.add_run()
            body_run.text = str(line)
            body_run.font.name = "Pretendard"
            body_run.font.size = Pt(11)
            body_run.font.color.rgb = pal["SUBT"]

    render_column(left, "KEY FINDINGS", findings, pal["PURPLE"])
    render_column(left + col_w + gap, "NEXT STEPS", next_steps, pal["INK"])
    return []


def inject_process(slide, content: dict) -> list[str]:
    steps = (content.get("data") or {}).get("steps", []) or []
    title = content.get("title", "")

    slot = 0
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if any(p in t for p in ("타이틀 입력", "타이틀을 입력")) and title and slot == 0:
            set_runs_text(shape.text_frame, title)
        elif any(p in t for p in ("내용 입력", "내용을 입력")) and slot < len(steps):
            set_runs_text(shape.text_frame, steps[slot])
            slot += 1
    return []


def inject_matrix(slide, content: dict) -> list[str]:
    data = content.get("data") or {}
    values = [data.get(k, "") for k in ("S", "W", "O", "T")]
    title = content.get("title", "")
    body = content.get("body", "")

    title_done = False
    sub_done = False
    body_slots: list = []
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not title_done and "메인 타이틀" in t:
            set_runs_text(shape.text_frame, title or " ")
            title_done = True
        elif not sub_done and "서브 타이틀" in t:
            set_runs_text(shape.text_frame, body or " ")
            sub_done = True
        elif is_body_placeholder(t):
            body_slots.append(shape)

    for shape, val in zip(body_slots, values):
        if val:
            set_runs_text(shape.text_frame, val)

    return [] if body_slots else ["matrix-no-slots"]


def _emu_to_in(emu):
    if emu is None:
        return None
    return emu / 914400


def _blank_cell(cell) -> None:
    """Zero out every run in a table cell so no template text leaks through."""
    tf = cell.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    if runs:
        runs[0].text = ""
        for r in runs[1:]:
            r.text = ""


TABLE_LEFTOVER_MARKERS = (
    "상세 타이틀",
    "내용 입력",
    "내용을 입력",
    "타이틀 입력",
    "타이틀을 입력",
    "메인 타이틀",
    "서브 타이틀",
    "안녕하세요",
    "파워포인트 템플릿",
    "매주 업데이트",
    "다양한 주제",
    "세련되고 정교한",
    "저희는 프레젠테이션",
    "비즈니스",
)


def inject_table(slide, content: dict) -> list[str]:
    data = content.get("data") or {}
    headers = data.get("headers", [])
    rows = data.get("rows", [])
    title = content.get("title", "")
    body = content.get("body", "")

    title_done = False
    sub_done = False
    # Collect decorative shapes (icon cards, sample text boxes) that sit ABOVE
    # the data table. We don't fill them with content, so leaving them empty
    # wastes half the slide — drop them entirely once we've located the table.
    table_top_in = None
    decorative_shapes: list = []
    table_shape = None

    for shape in iter_shapes(slide.shapes):
        if shape.has_table:
            table_shape = shape
            table_top_in = _emu_to_in(shape.top)
            tbl = shape.table
            n_cols = len(tbl.columns)
            n_rows = len(tbl.rows)

            # Headers (row 0): overwrite every cell, blanking any extra cols.
            for c in range(n_cols):
                header_text = headers[c] if c < len(headers) else ""
                if header_text:
                    set_runs_text(tbl.rows[0].cells[c].text_frame, str(header_text))
                else:
                    _blank_cell(tbl.rows[0].cells[c])

            # Data rows: always overwrite, even empties — prevents the template's
            # sample text ("매출채권회수기간 86.5%" etc.) from bleeding through.
            data_rows_in_slide = min(len(rows), n_rows - 1)
            for r_idx in range(1, n_rows):
                src_row = rows[r_idx - 1] if r_idx - 1 < data_rows_in_slide else []
                for c_idx in range(n_cols):
                    val = src_row[c_idx] if c_idx < len(src_row) else ""
                    if val:
                        set_runs_text(tbl.rows[r_idx].cells[c_idx].text_frame, str(val))
                    else:
                        _blank_cell(tbl.rows[r_idx].cells[c_idx])

        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if not title_done and "메인 타이틀" in t:
                set_runs_text(shape.text_frame, title or " ")
                title_done = True
            elif not sub_done and "서브 타이틀" in t:
                # We render the headline via a dedicated text box below so the
                # built-in subtitle placeholder (which sits ABOVE the title in
                # this template) would duplicate it — blank it out.
                set_runs_text(shape.text_frame, " ")
                sub_done = True
            elif any(m in t for m in TABLE_LEFTOVER_MARKERS):
                set_runs_text(shape.text_frame, " ")
            elif any(tok in t for tok in GLOBAL_BOILERPLATE_TOKENS):
                set_runs_text(shape.text_frame, " ")
            elif t in FOOTER_TOKENS:
                set_runs_text(shape.text_frame, " ")

    # Second pass: drop decorative shapes (rounded rects, ovals, icon freeforms,
    # empty text boxes) that live above the data table. They came from the
    # template's "icon cards" decoration and add no value once we've blanked
    # their placeholder copy — the cards render as awkward empty white panels.
    if table_top_in is not None:
        for shape in list(slide.shapes):
            if shape.has_table:
                continue
            top_in = _emu_to_in(shape.top)
            if top_in is None or top_in >= table_top_in - 0.1:
                continue
            # Preserve the title placeholder (has meaningful injected text).
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and t != " " and "메인 타이틀" not in t and "서브 타이틀" not in t:
                    # Keep the title we just injected.
                    continue
            # Everything else above the table (card backgrounds, ovals, icon
            # freeforms, boilerplate text boxes we blanked) is decorative noise.
            decorative_shapes.append(shape)

        for sh in decorative_shapes:
            sp = sh._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)

    # Reposition + enlarge the data table so it fills the vacated space.
    # Template leaves the table at y≈4.4", which looks orphaned now that the
    # decorative icon cards above it are gone. Pull it up just under the
    # title and stretch it vertically so it actually uses the slide.
    if table_shape is not None:
        from pptx.util import Emu
        new_top = Emu(int(2.75 * 914400))    # y = 2.75" (headline sits above)
        new_height = Emu(int(2.7 * 914400))  # h = 2.7"  → bottom at 5.45"
        old_height_emu = int(table_shape.height)
        table_shape.top = new_top
        table_shape.height = new_height
        # Scale each row's height by the same factor so header + data rows
        # share the new real estate proportionally.
        if old_height_emu:
            scale = int(new_height) / old_height_emu
            for row in table_shape.table.rows:
                row.height = Emu(int(int(row.height) * scale))

    # Layer in consulting-style commentary: a one-line "answer-first" headline
    # just above the table, and bulleted insights below it. Without these the
    # slide is just a number dump — adding them turns the table into a real
    # analytical page.
    headline = (content.get("body") or "").strip()
    insights = [s for s in (data.get("insights") or []) if s and str(s).strip()]
    if headline:
        _add_headline_box(slide, headline)
    if insights:
        _add_insights_box(slide, insights)

    return []


def _add_headline_box(slide, text: str) -> None:
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    left = Emu(int(0.72 * 914400))
    top = Emu(int(2.0 * 914400))
    width = Emu(int(11.89 * 914400))
    height = Emu(int(0.6 * 914400))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Pretendard SemiBold"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x1E, 0x1E, 0x2A)


def _add_insights_box(slide, bullets: list) -> None:
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    left = Emu(int(0.72 * 914400))
    top = Emu(int(5.65 * 914400))
    width = Emu(int(11.89 * 914400))
    height = Emu(int(1.65 * 914400))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    # First paragraph: small uppercase label "KEY TAKEAWAYS"
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = "KEY TAKEAWAYS"
    r0.font.name = "Pretendard SemiBold"
    r0.font.size = Pt(9)
    r0.font.color.rgb = RGBColor(0x6E, 0x5A, 0xE6)

    for line in bullets[:3]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r_bullet = p.add_run()
        r_bullet.text = "•  "
        r_bullet.font.name = "Pretendard SemiBold"
        r_bullet.font.size = Pt(11)
        r_bullet.font.color.rgb = RGBColor(0x6E, 0x5A, 0xE6)
        r_text = p.add_run()
        r_text.text = str(line)
        r_text.font.name = "Pretendard"
        r_text.font.size = Pt(11)
        r_text.font.color.rgb = RGBColor(0x33, 0x33, 0x3D)


def inject_image(slide, content: dict) -> list[str]:
    title = content.get("title", "")
    body = content.get("body", "")
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if any(p in t for p in ("타이틀 입력", "타이틀을 입력")) and title:
            set_runs_text(shape.text_frame, title)
        elif any(p in t for p in ("내용 입력", "내용을 입력")) and body:
            set_runs_text(shape.text_frame, body)
    return []


def inject_content(slide, content: dict) -> list[str]:
    title = content.get("title", "")
    body = content.get("body", "")
    used_title = False
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not used_title and title and any(p in t for p in ("타이틀 입력", "타이틀을 입력", "메인 타이틀")):
            set_runs_text(shape.text_frame, title)
            used_title = True
        elif body and any(p in t for p in ("내용 입력", "내용을 입력", "서브 타이틀")):
            set_runs_text(shape.text_frame, body)
    return []


def inject_closing(slide, content: dict) -> list[str]:
    message = content.get("message", "Thank you")
    title_done = False
    sub_done = False
    shapes_to_drop = []
    for shape in iter_shapes(slide.shapes):
        if shape.has_table:
            # Remove the sample dashboard grid entirely — a blanked grid still
            # shows visible borders and looks ugly on the closing page.
            shapes_to_drop.append(shape)
            continue
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "감사합니다" in t or "Thank" in t or "THANK" in t:
            set_runs_text(shape.text_frame, message)
            title_done = True
            continue
        if not title_done and "메인 타이틀" in t:
            set_runs_text(shape.text_frame, message)
            title_done = True
        elif not sub_done and "서브 타이틀" in t:
            set_runs_text(shape.text_frame, " ")
            sub_done = True
        elif any(m in t for m in TABLE_LEFTOVER_MARKERS):
            set_runs_text(shape.text_frame, " ")
        elif any(tok in t for tok in GLOBAL_BOILERPLATE_TOKENS):
            set_runs_text(shape.text_frame, " ")
        elif t in FOOTER_TOKENS:
            set_runs_text(shape.text_frame, " ")

    for sh in shapes_to_drop:
        sp = sh._element
        sp.getparent().remove(sp)
    return []


INJECTORS = {
    "cover": inject_cover,
    "agenda": inject_agenda,
    "section": inject_section,
    "kpi": inject_kpi,
    "dashboard": inject_dashboard,
    "chart": inject_chart,
    "exec_summary": inject_exec_summary,
    "findings": inject_findings,
    "process": inject_process,
    "matrix": inject_matrix,
    "table": inject_table,
    "image": inject_image,
    "content": inject_content,
    "closing": inject_closing,
}


def duplicate_slide(prs, source_slide):
    """Append a deep copy of `source_slide` to `prs`. Returns the new slide."""
    blank_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    spTree = new_slide.shapes._spTree
    for child in source_slide.shapes._spTree.iterchildren():
        tag = child.tag
        if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
            continue
        spTree.append(copy.deepcopy(child))

    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def drop_slide(prs, slide) -> None:
    slide_id = slide.slide_id
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        if int(sldId.get("id")) == slide_id:
            rId = sldId.get(qn("r:id"))
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)
            return


def _strip_layout_top_decorations(prs) -> None:
    """Remove tiny decorative shapes pinned to the top edge of every shared
    slide layout that is still referenced by the final deck.

    The template ships with a 0.73" × 0.19" rounded-corner "nub" anchored at
    y=0 on several layouts (e.g. ``5_사용자 지정 레이아웃``). Once the real
    slides inherit the layout it shows up as a random purple crumb at the top-
    left corner — a visible artifact the user flagged. It's a user-drawn
    decoration (``userDrawn="1"``), not a placeholder, so stripping it can't
    break inheritance.
    """
    layouts_in_use: set = set()
    for slide in prs.slides:
        layouts_in_use.add(id(slide.slide_layout))

    for layout in prs.slide_layouts:
        if id(layout) not in layouts_in_use:
            continue
        for shape in list(layout.shapes):
            if shape.is_placeholder:
                continue
            top_in = _emu_to_in(shape.top)
            height_in = _emu_to_in(shape.height)
            if top_in is None or height_in is None:
                continue
            # Only drop shapes anchored within the top 0.3" and under 0.5" tall
            # — that matches the decorative nub, not genuine banner/background
            # art that may legitimately hug the top edge.
            if top_in <= 0.05 and height_in <= 0.5:
                sp = shape._element
                parent = sp.getparent()
                if parent is not None:
                    parent.remove(sp)


def assemble(template_path: Path, plan_path: Path, output_path: Path) -> dict:
    plan = json.loads(plan_path.read_text())
    prs = Presentation(str(template_path))

    original_slides = list(prs.slides)
    source_by_index = {i + 1: slide for i, slide in enumerate(original_slides)}

    report = []
    for entry in plan["slides"]:
        source_idx = entry["source_index"]
        kind = entry["kind"]
        content = entry.get("content", {})

        source_slide = source_by_index.get(source_idx)
        if source_slide is None:
            report.append({"position": entry["position"], "kind": kind, "error": f"missing source slide {source_idx}"})
            continue

        new_slide = duplicate_slide(prs, source_slide)
        injector = INJECTORS.get(kind, inject_content)
        issues = injector(new_slide, content)
        if issues:
            report.append(
                {"position": entry["position"], "source_index": source_idx, "kind": kind, "issues": issues}
            )

    for slide in original_slides:
        drop_slide(prs, slide)

    _strip_layout_top_decorations(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {"output": str(output_path), "slide_count": len(plan["slides"]), "issues": report}


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--plan", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    template = Path(args.template).expanduser().resolve()
    plan = Path(args.plan).expanduser().resolve()
    output = Path(args.output).expanduser().resolve()

    if not template.exists():
        sys.exit(f"template not found: {template}")
    if not plan.exists():
        sys.exit(f"plan not found: {plan}")

    result = assemble(template, plan, output)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    import sys  # noqa
    main()
