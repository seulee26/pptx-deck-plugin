"""Classify each slide of a PPTX template by visual/content type and emit a
classification.json catalog that slide-matcher consumes.

Usage:
    python3 classify_template.py <template.pptx> <classification.json>

Heuristics match the slide kinds recognised by the rest of the pipeline:
cover / agenda / section / kpi / process / matrix / table / image / content / closing.
"""

from __future__ import annotations

import json
import re
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BIG_NUM_RE = re.compile(r"\d{1,3}(?:,\d{3})+|\d{1,3}%|\b\d{3,}\b|\b\d{2}%")
STEP_RE = re.compile(r"(?<!\d)0[1-9](?!\d)")
SWOT_LETTERS = {"S", "W", "O", "T"}


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def table_text(shape) -> str:
    if not shape.has_table:
        return ""
    parts = []
    for row in shape.table.rows:
        for cell in row.cells:
            t = cell.text_frame.text.strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


def analyze_slide(slide) -> dict:
    has_chart = has_table = has_picture = False
    texts: list[str] = []
    shape_count = 0
    table_cells = 0
    for shape in iter_shapes(slide.shapes):
        shape_count += 1
        if shape.has_chart:
            has_chart = True
        if shape.has_table:
            has_table = True
            table_cells += len(shape.table.rows) * len(shape.table.columns)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True
        t = shape_text(shape)
        if t:
            texts.append(t)
        t = table_text(shape)
        if t:
            texts.append(t)
    joined = "\n".join(texts)
    return {
        "has_chart": has_chart,
        "has_table": has_table,
        "has_picture": has_picture,
        "shape_count": shape_count,
        "table_cells": table_cells,
        "text": joined,
        "text_len": len(joined.replace(" ", "")),
    }


def classify(info: dict, index: int, total: int) -> str:
    text = info["text"]
    lower = text.lower()

    if index == 1:
        return "cover"
    if index == total:
        return "closing"
    if "thank" in lower or "q&a" in lower or "감사합니다" in text:
        return "closing"

    step_matches = sorted(set(STEP_RE.findall(text)))
    if (
        "PRESENTATION TEMPLATE" in text
        and "LOGO HERE" in text
        and len(step_matches) >= 1
        and info["text_len"] < 120
    ):
        return "section"

    tokens = set(re.findall(r"\b[A-Z]\b", text))
    if (
        any(k in text for k in ("Strength", "Weakness", "Opportunity", "Threat"))
        or SWOT_LETTERS.issubset(tokens)
    ):
        return "matrix"

    if any(k in lower for k in ("table of content", "agenda", "contents")) or "목차" in text:
        return "agenda"
    if len(step_matches) >= 5 and info["text_len"] < 260 and info["shape_count"] <= 25:
        return "agenda"

    if info["has_table"] and info["table_cells"] >= 9 and not info["has_picture"]:
        return "table"

    if len(step_matches) >= 3:
        return "process"

    if BIG_NUM_RE.search(text):
        return "kpi"

    if text.count("타이틀 입력") >= 3 and info["text_len"] < 260 and info["shape_count"] <= 20:
        return "org"

    if info["has_picture"] or (info["has_table"] and info["text_len"] < 260):
        return "image"

    return "content"


def main() -> None:
    if len(sys.argv) != 3:
        sys.exit("usage: classify_template.py <template.pptx> <classification.json>")

    src = Path(sys.argv[1]).expanduser().resolve()
    dst = Path(sys.argv[2]).expanduser().resolve()
    if not src.exists():
        sys.exit(f"template not found: {src}")

    prs = Presentation(str(src))
    total = len(prs.slides)
    records = []
    for i, slide in enumerate(prs.slides, start=1):
        info = analyze_slide(slide)
        kind = classify(info, i, total)
        records.append(
            {
                "index": i,
                "kind": kind,
                "has_chart": info["has_chart"],
                "has_table": info["has_table"],
                "has_picture": info["has_picture"],
                "shape_count": info["shape_count"],
                "text_len": info["text_len"],
                "text_preview": info["text"][:180].replace("\n", " | "),
            }
        )

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(
        json.dumps({"total": total, "slides": records}, ensure_ascii=False, indent=2)
    )

    counts = Counter(r["kind"] for r in records)
    print(f"classified {total} slides → {dst}")
    for k, n in counts.most_common():
        print(f"  {k}: {n}")


if __name__ == "__main__":
    main()
